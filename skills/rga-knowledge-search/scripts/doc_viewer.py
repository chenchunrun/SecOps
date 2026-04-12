#!/usr/bin/env python3
"""
统一文档内容查看器 - 支持 PDF、Office、电子书、压缩包等格式
"""

import sys
import os
import argparse
import tempfile
import shutil
from pathlib import Path

# 默认最大输出字符数 (约 12500 tokens)
DEFAULT_MAX_CHARS = 50000


def check_dependencies():
    """检查并报告依赖状态"""
    deps = {
        "pdfplumber": {"installed": False, "purpose": "PDF 文本提取"},
        "python-docx": {"installed": False, "purpose": "Word 文档"},
        "openpyxl": {"installed": False, "purpose": "Excel 文件"},
        "python-pptx": {"installed": False, "purpose": "PowerPoint 文件"},
        "ebooklib": {"installed": False, "purpose": "EPUB 电子书"},
        "chardet": {"installed": False, "purpose": "编码检测"},
    }

    try:
        import pdfplumber
        deps["pdfplumber"]["installed"] = True
    except ImportError:
        pass

    try:
        import docx
        deps["python-docx"]["installed"] = True
    except ImportError:
        pass

    try:
        import openpyxl
        deps["openpyxl"]["installed"] = True
    except ImportError:
        pass

    try:
        import pptx
        deps["python-pptx"]["installed"] = True
    except ImportError:
        pass

    try:
        import ebooklib
        deps["ebooklib"]["installed"] = True
    except ImportError:
        pass

    try:
        import chardet
        deps["chardet"]["installed"] = True
    except ImportError:
        pass

    return deps


def extract_pdf(file_path: str, pages: list = None) -> str:
    """提取 PDF 文本内容"""
    try:
        import pdfplumber
    except ImportError:
        return "[错误] 需要安装 pdfplumber: pip install pdfplumber"

    text_parts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            target_pages = pages if pages else range(total_pages)

            for i in target_pages:
                if i >= total_pages:
                    continue
                page = pdf.pages[i]
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(f"--- 第 {i+1}/{total_pages} 页 ---\n{text}")

            if not text_parts:
                return "[提示] PDF 无可提取文本，可能是扫描件，建议使用 OCR"

            return "\n\n".join(text_parts)
    except Exception as e:
        return f"[错误] PDF 读取失败: {e}"


def extract_docx(file_path: str) -> str:
    """提取 Word 文档内容"""
    try:
        from docx import Document
    except ImportError:
        return "[错误] 需要安装 python-docx: pip install python-docx"

    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # 提取表格
        tables_text = []
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                tables_text.append(f"[表格 {i+1}]\n" + "\n".join(rows))

        result = "\n\n".join(paragraphs)
        if tables_text:
            result += "\n\n--- 表格内容 ---\n" + "\n\n".join(tables_text)

        return result if result.strip() else "[提示] 文档为空"
    except Exception as e:
        return f"[错误] Word 文档读取失败: {e}"


def extract_xlsx(file_path: str, sheet_name: str = None) -> str:
    """提取 Excel 文件内容"""
    try:
        import openpyxl
    except ImportError:
        return "[错误] 需要安装 openpyxl: pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames

        result_parts = []
        for sheet in sheets:
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cell_values = [str(c) if c is not None else "" for c in row]
                if any(cell_values):
                    rows.append(" | ".join(cell_values))

            if rows:
                result_parts.append(f"--- Sheet: {sheet} ---\n" + "\n".join(rows))

        return "\n\n".join(result_parts) if result_parts else "[提示] 工作簿为空"
    except Exception as e:
        return f"[错误] Excel 读取失败: {e}"


def extract_pptx(file_path: str) -> str:
    """提取 PowerPoint 内容"""
    try:
        from pptx import Presentation
    except ImportError:
        return "[错误] 需要安装 python-pptx: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides_text.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(texts))

        return "\n\n".join(slides_text) if slides_text else "[提示] 演示文稿为空"
    except Exception as e:
        return f"[错误] PowerPoint 读取失败: {e}"


def extract_epub(file_path: str) -> str:
    """提取 EPUB 电子书内容"""
    try:
        import ebooklib
        from ebooklib import epub
        from html.parser import HTMLParser
    except ImportError:
        return "[错误] 需要安装 ebooklib: pip install ebooklib"

    class HTMLTextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []

        def handle_data(self, data):
            self.text.append(data)

        def get_text(self):
            return " ".join(self.text)

    try:
        book = epub.read_epub(file_path)
        chapters = []

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                parser = HTMLTextExtractor()
                content = item.get_content().decode("utf-8", errors="ignore")
                parser.feed(content)
                text = parser.get_text().strip()
                if text:
                    chapters.append(text)

        return "\n\n---\n\n".join(chapters) if chapters else "[提示] 电子书为空"
    except Exception as e:
        return f"[错误] EPUB 读取失败: {e}"


def extract_archive(file_path: str, target_file: str = None) -> str:
    """列出或提取压缩包内容"""
    import zipfile
    import tarfile

    ext = Path(file_path).suffix.lower()

    try:
        if ext == ".zip":
            with zipfile.ZipFile(file_path, "r") as zf:
                if target_file:
                    try:
                        content = zf.read(target_file)
                        try:
                            return content.decode("utf-8")
                        except UnicodeDecodeError:
                            try:
                                import chardet
                                detected = chardet.detect(content)
                                return content.decode(detected["encoding"] or "utf-8", errors="replace")
                            except ImportError:
                                return content.decode("utf-8", errors="replace")
                    except KeyError:
                        return f"[错误] 文件 '{target_file}' 不在压缩包中"
                else:
                    files = zf.namelist()
                    return f"压缩包内容 ({len(files)} 个文件):\n" + "\n".join(f"  {f}" for f in files)

        elif ext in [".tar", ".gz", ".tgz", ".bz2", ".xz"]:
            mode = "r:*"
            with tarfile.open(file_path, mode) as tf:
                if target_file:
                    try:
                        member = tf.getmember(target_file)
                        f = tf.extractfile(member)
                        if f:
                            content = f.read()
                            try:
                                return content.decode("utf-8")
                            except UnicodeDecodeError:
                                return content.decode("utf-8", errors="replace")
                        else:
                            return "[错误] 无法读取该文件（可能是目录）"
                    except KeyError:
                        return f"[错误] 文件 '{target_file}' 不在压缩包中"
                else:
                    files = tf.getnames()
                    return f"压缩包内容 ({len(files)} 个文件):\n" + "\n".join(f"  {f}" for f in files)
        else:
            return f"[错误] 不支持的压缩格式: {ext}"

    except Exception as e:
        return f"[错误] 压缩包读取失败: {e}"


def extract_text(file_path: str) -> str:
    """提取纯文本文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            import chardet
            with open(file_path, "rb") as f:
                raw = f.read()
            detected = chardet.detect(raw)
            return raw.decode(detected["encoding"] or "utf-8", errors="replace")
        except ImportError:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception as e:
        return f"[错误] 文本读取失败: {e}"


def view_document(file_path: str, **kwargs) -> str:
    """根据文件类型自动选择提取方法"""
    path = Path(file_path)

    if not path.exists():
        return f"[错误] 文件不存在: {file_path}"

    ext = path.suffix.lower()

    # 处理复合扩展名
    if path.name.endswith(".tar.gz") or path.name.endswith(".tgz"):
        ext = ".tar.gz"
    elif path.name.endswith(".tar.bz2"):
        ext = ".tar.bz2"
    elif path.name.endswith(".tar.xz"):
        ext = ".tar.xz"

    extractors = {
        ".pdf": lambda: extract_pdf(file_path, kwargs.get("pages")),
        ".docx": lambda: extract_docx(file_path),
        ".doc": lambda: "[提示] .doc 格式需要 antiword 或转换为 .docx",
        ".xlsx": lambda: extract_xlsx(file_path, kwargs.get("sheet")),
        ".xls": lambda: "[提示] .xls 格式建议转换为 .xlsx",
        ".pptx": lambda: extract_pptx(file_path),
        ".ppt": lambda: "[提示] .ppt 格式建议转换为 .pptx",
        ".epub": lambda: extract_epub(file_path),
        ".zip": lambda: extract_archive(file_path, kwargs.get("target")),
        ".tar": lambda: extract_archive(file_path, kwargs.get("target")),
        ".tar.gz": lambda: extract_archive(file_path, kwargs.get("target")),
        ".tgz": lambda: extract_archive(file_path, kwargs.get("target")),
        ".tar.bz2": lambda: extract_archive(file_path, kwargs.get("target")),
        ".tar.xz": lambda: extract_archive(file_path, kwargs.get("target")),
        ".gz": lambda: extract_archive(file_path, kwargs.get("target")),
        ".txt": lambda: extract_text(file_path),
        ".md": lambda: extract_text(file_path),
        ".json": lambda: extract_text(file_path),
        ".xml": lambda: extract_text(file_path),
        ".csv": lambda: extract_text(file_path),
        ".log": lambda: extract_text(file_path),
    }

    if ext in extractors:
        return extractors[ext]()
    else:
        # 尝试作为文本读取
        try:
            return extract_text(file_path)
        except Exception:
            return f"[错误] 不支持的文件格式: {ext}"


def main():
    parser = argparse.ArgumentParser(
        description="统一文档内容查看器",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  %(prog)s document.pdf                    # 查看 PDF
  %(prog)s document.pdf --pages 0,1,2      # 查看前3页
  %(prog)s report.docx                     # 查看 Word 文档
  %(prog)s data.xlsx                       # 查看 Excel
  %(prog)s data.xlsx --sheet "Sheet1"      # 查看指定工作表
  %(prog)s slides.pptx                     # 查看 PowerPoint
  %(prog)s book.epub                       # 查看电子书
  %(prog)s archive.zip                     # 列出压缩包内容
  %(prog)s archive.zip --target readme.txt # 查看压缩包内文件
  %(prog)s --check-deps                    # 检查依赖状态
        """
    )

    parser.add_argument("file", nargs="?", help="要查看的文件路径")
    parser.add_argument("--pages", help="PDF 页码，逗号分隔 (从0开始)")
    parser.add_argument("--sheet", help="Excel 工作表名称")
    parser.add_argument("--target", help="压缩包内的目标文件")
    parser.add_argument("--check-deps", action="store_true", help="检查依赖状态")
    parser.add_argument("--max-lines", type=int, default=0, help="限制输出行数 (0=不限制)")
    parser.add_argument("--max-chars", "-c", type=int, default=DEFAULT_MAX_CHARS,
                        help=f"最大输出字符数 (默认: {DEFAULT_MAX_CHARS}, 0=不限制)")
    parser.add_argument("--no-limit", action="store_true", help="取消所有输出限制")

    args = parser.parse_args()

    if args.check_deps:
        deps = check_dependencies()
        print("依赖状态检查:\n")
        for name, info in deps.items():
            status = "[+] 已安装" if info["installed"] else "[-] 未安装"
            print(f"  {name}: {status} ({info['purpose']})")

        missing = [n for n, i in deps.items() if not i["installed"]]
        if missing:
            print(f"\n安装缺失依赖: pip install {' '.join(missing)}")
        return

    if not args.file:
        parser.print_help()
        return

    # 解析页码
    pages = None
    if args.pages:
        try:
            pages = [int(p.strip()) for p in args.pages.split(",")]
        except ValueError:
            print("[错误] 页码格式错误，应为逗号分隔的数字")
            return

    result = view_document(
        args.file,
        pages=pages,
        sheet=args.sheet,
        target=args.target
    )

    original_chars = len(result)
    original_lines = result.count('\n') + 1
    truncated = False

    # 应用限制
    if not args.no_limit:
        # 先按行数截断
        if args.max_lines > 0:
            lines = result.split("\n")
            if len(lines) > args.max_lines:
                result = "\n".join(lines[:args.max_lines])
                truncated = True

        # 再按字符数截断
        if args.max_chars > 0 and len(result) > args.max_chars:
            result = result[:args.max_chars]
            # 尝试在最后一个完整行处截断
            last_newline = result.rfind('\n')
            if last_newline > args.max_chars * 0.8:
                result = result[:last_newline]
            truncated = True

    print(result)

    if truncated:
        final_chars = len(result)
        final_lines = result.count('\n') + 1
        print(f"\n[已截断] 输出 {final_chars}/{original_chars} 字符, {final_lines}/{original_lines} 行",
              file=sys.stderr)
        print(f"[提示] 使用 --max-chars N 调整限制，或 --no-limit 取消限制", file=sys.stderr)


if __name__ == "__main__":
    main()
